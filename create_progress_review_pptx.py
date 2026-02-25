"""Generate EduConnect Project Progress Review PowerPoint"""
from pptx import Presentation
from pptx.util import Inches, Pt

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, bullets, subbullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.add_paragraph() if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(14)
        if subbullets and i < len(subbullets) and subbullets[i]:
            for sub in subbullets[i]:
                sp = body.add_paragraph()
                sp.text = sub
                sp.level = 1
                sp.font.size = Pt(12)

def add_para(prs, title, paras):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for ptext in paras:
        p = body.add_paragraph()
        p.text = ptext
        p.font.size = Pt(14)
        p.space_after = Pt(6)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Title
add_title_slide(prs, "EduConnect", "Project Progress Review Presentation")
add_title_slide(prs, "Project Overview", "A study partner platform connecting students, tracking progress, and providing learning resources")

# 1. Progress Against Previous Commitments
add_content_slide(prs, "1. Progress Against Previous Commitments", [
    "Deliverables committed & status:",
    "User Auth & Registration – Completed",
    "Recommendation Engine – Completed",
    "Study Tracking & Analytics – Completed",
    "Learning Resources & Quizzes – Completed",
    "Study Groups – Completed",
    "BERTopic NLP / Topic Modeling – Completed",
    "Dashboard & Responsive Analytics – Completed",
    "Feedback System – Completed",
], [[
    "Multi-step signup, login, profile, localStorage",
    "Weighted Jaccard matching, match scores",
    "Session timer, weekly hours, StudyAnalytics page",
    "Curated resources, field quizzes, final tests",
    "Interest-based groups, 30% minimum match",
    "Background topic extraction from feedback/chat",
    "Charts, empty states, mobile-responsive",
    "Submit feedback, history, status tracking",
]])

add_para(prs, "1. Progress Against Previous Commitments (cont.)", [
    "b) Milestone completion estimate: ~85–90%",
    "",
    "Justification: Core platform is fully functional. All major features are implemented: auth, recommendations, study tracking, resources with quizzes, study groups, NLP topic modeling, and analytics. Remaining work is polish, backend migration, and optional enhancements (chat, push notifications).",
])

# 2. Technical Implementation
add_content_slide(prs, "2. Technical Implementation Since Last Review", [
    "Components implemented:",
    "Frontend: React 18, Vite, React Router, Framer Motion, Recharts",
    "Recommendation Engine: Weighted Jaccard similarity (CS interests 40%, skills, learning style, etc.)",
    "Study tracking: Session timer, 5-min inactivity pause, manual stop, recordManualStudySession",
    "Quizzes: Field-based quizzes with proficiency (Beginner/Intermediate/Advanced/Expert)",
    "BERTopic: Flask backend for topic modeling on feedback & chat",
    "Responsive analytics: Charts with ResponsiveContainer, empty states, animations",
], [[
    "Tech stack and architecture",
    "Matching algorithm in recommendationEngine.js",
    "StudyContext + AuthContext integration",
    "quizData.js, Quiz.jsx, recordFieldProgress",
    "backend/app.py, nlpBackgroundService.js",
    "StudyAnalytics.css media queries, Y-axis domains",
]])

add_para(prs, "2. Technical Iterations & Evidence", [
    "Key iteration: NLP Insights moved to background (no user click needed). Topics auto-analyzed on login, cached 30 min.",
    "",
    "Evidence: Core functionality works end-to-end. Recommendation matching produces ranked partners. Study sessions log hours. Quizzes update field progress. Topic modeling runs via backend. Charts display weekly hours and progress. Documentation in DOCUMENTATION.md.",
])

# 3. Current Working State
add_content_slide(prs, "3. Current Working State of the Project", [
    "a) Current prototype: Full React SPA with Flask backend for NLP. Runs locally (npm run dev + python app.py).",
    "b) Live-demonstrable functionality:",
    "c) Incomplete / under development:",
    "d) Biggest technical challenges:",
], [[
    "",
    "Sign up, login, dashboard, recommendations, study groups, resources, quizzes, analytics, feedback, topic display on dashboard.",
    "Backend auth (still localStorage), real-time chat, push notifications, production deployment.",
    "Client-side only auth; scaling BERTopic for larger corpus.",
]])

# 4. Response to Previous Feedback
add_para(prs, "4. Response to Previous Feedback", [
    "a) If no prior review feedback was received, state: 'First progress review – no prior feedback to address.'",
    "",
    "b) If feedback was received, list each point and how it was addressed.",
    "",
    "c) Any feedback not implemented: Justify (e.g., deferred to later phase, scope change).",
])

# 5. Individual Contribution
add_content_slide(prs, "5. Individual Contribution", [
    "Template – customize per team member:",
    "Member 1: Name – Modules: [e.g., Auth, Recommendation Engine] – Work: [specific tasks] – Evidence: [files, commits]",
    "Member 2: Name – Modules: [e.g., Study Tracking, Analytics] – Work: [specific tasks] – Evidence: [files, commits]",
    "Member 3: Name – Modules: [e.g., Resources, Quizzes, NLP] – Work: [specific tasks] – Evidence: [files, commits]",
], [])

# 6. Updated Plan and Risk Assessment
add_content_slide(prs, "6. Updated Plan and Risk Assessment", [
    "a) Remaining tasks: Backend auth, deployment, optional chat/notifications, UI polish",
    "b) Timeline: [Specify target completion date]",
    "c) Risks & mitigations:",
    "d) On track: Yes / No – [brief explanation]",
], [[
    "",
    "",
    "Risk: localStorage limits scalability → Mitigation: Migrate to backend DB. Risk: BERTopic slow on large data → Mitigation: Batch processing, cache topics.",
    "",
]])

add_para(prs, "6. Risks & Mitigations (detail)", [
    "Technical risk: Client-side storage → Mitigation: Plan backend API + database (e.g., PostgreSQL).",
    "Technical risk: NLP performance → Mitigation: Run topic modeling asynchronously; cache results.",
    "Project risk: Scope creep → Mitigation: Stick to MVP; defer chat/notifications to post-MVP.",
])

# Summary
add_title_slide(prs, "Summary", "EduConnect: Core platform complete. Auth, recommendations, study tracking, resources, quizzes, study groups, NLP, and analytics implemented. On track for completion.")

out_path = r"c:\Users\USER\OneDrive\Desktop\The Actual Educonnect (2)\EduConnect_Project_Progress_Review.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
